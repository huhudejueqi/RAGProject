"""LangChain 文档加载器注册表。新增文件类型只需维护注册项，无需 if/elif 分支。"""

from __future__ import annotations

import importlib
import importlib.util
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

import fitz
from langchain_core.documents import Document
from docx import Document as DocxDocument
from pptx import Presentation

from qa_core.config.settings import get_settings
from qa_core.config.logging_config import get_logger
from qa_core.indexing.ocr_review import parse_ocr_review_metadata
from qa_core.indexing.table_documents import load_table_file
logger = get_logger(__name__)

DOCLING_SUFFIXES = {".pdf", ".docx", ".pptx", ".html", ".htm"}

class SupportsLoad(Protocol):
    """文档加载器最小协议。不直接导入 LangChain BaseLoader，避免提前拉起 heavy 依赖。

    调用顺序：入库脚本或索引服务 -> SupportsLoad。
    """
    def load(self) -> list[Document]:
        """加载文件并返回 LangChain Document 列表。

        调用顺序：入库脚本或索引服务 -> SupportsLoad.load()。
        """
        ...

LoaderFactory = Callable[[Path], SupportsLoad]

@dataclass(frozen=True)
class DocumentLoaderSpec:
    """文档加载器注册项。保存 factory 而非裸 loader class，因不同 loader 构造参数不同。

    调用顺序：入库脚本或索引服务 -> DocumentLoaderSpec。
    """

    suffixes: tuple[str, ...]
    factory: LoaderFactory
    description: str

    def create_loader(self, path: Path) -> SupportsLoad:
        """创建实际 loader 实例。

        工厂函数统一接收 Path，内部自行决定是否需要 encoding、mode 等参数。这样
        `load_file()` 不需要知道 TextLoader、PyPDFLoader、Docx2txtLoader 的构造差异。

        调用顺序：入库脚本或索引服务 -> DocumentLoaderSpec.create_loader()。
        """
        return self.factory(path)

def _utf8_text_loader(path: Path) -> SupportsLoad:
    """创建 UTF-8 文本加载器。.md 也走 TextLoader 以保留 Markdown 标题给后续结构化切分。

    调用顺序：入库脚本或索引服务 -> _utf8_text_loader()。
    """
    return ProjectTextLoader(path)

def _pdf_loader(path: Path) -> SupportsLoad:
    """创建 PDF 文本层加载器。使用 PyMuPDF，适合有中文文本层的业务 PDF。

    调用顺序：入库脚本或索引服务 -> _pdf_loader()。
    """
    return PyMuPdfLoader(path)

def _word_loader(path: Path) -> SupportsLoad:
    """创建 Word 文档加载器。优先使用当前依赖里的 python-docx，避免额外依赖 docx2txt。

    调用顺序：入库脚本或索引服务 -> _word_loader()。
    """
    return PythonDocxLoader(path)


def _powerpoint_loader(path: Path) -> SupportsLoad:
    """创建 PPT/PPTX 文档加载器。优先使用当前依赖里的 python-pptx。

    调用顺序：入库脚本或索引服务 -> _powerpoint_loader()。
    """
    return PythonPptxLoader(path)


class DoclingLoader:
    """Docling 增强 loader，把复杂版面资料转换成 Markdown 后进入统一入库链路。

    调用顺序：入库脚本或索引服务 -> DoclingLoader。
    """

    def __init__(self, path: Path) -> None:
        """初始化 Docling 文档加载器。

        参数：
            path: 待解析的文档文件路径。

        调用顺序：入库脚本或索引服务 -> DoclingLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """使用 Docling 将复杂版面文档转为 Markdown 格式的 LangChain Document。

        调用顺序：入库脚本或索引服务 -> DoclingLoader.load()。
        """
        if (
            importlib.util.find_spec("docling") is None
            or importlib.util.find_spec("docling.document_converter") is None
        ):
            raise RuntimeError("Docling 解析后端不可用：未安装 docling。")

        converter_module = importlib.import_module("docling.document_converter")
        document_converter = converter_module.DocumentConverter()
        result = document_converter.convert(str(self.path))
        content = result.document.export_to_markdown().strip()
        metadata = {
            "file_type": self.path.suffix.lower(),
            "parser_backend": "docling",
            "docling_format": "markdown",
        }
        return [Document(page_content=content, metadata=metadata)] if content else []


def _docling_loader(path: Path) -> SupportsLoad:
    """创建 Docling 增强解析 loader。

    调用顺序：入库脚本或索引服务 -> _docling_loader()。
    """
    return DoclingLoader(path)


def _use_docling_for(path: Path) -> bool:
    """判断当前文件是否启用 Docling。表格文件保持项目行级 loader，避免丢失行列语义。

    调用顺序：入库脚本或索引服务 -> _use_docling_for()。
    """
    return get_settings().document_parser_backend == "docling" and path.suffix.lower() in DOCLING_SUFFIXES


class PythonDocxLoader:
    """轻量 DOCX loader，读取正文段落和表格文本。

    调用顺序：入库脚本或索引服务 -> PythonDocxLoader。
    """

    def __init__(self, path: Path) -> None:
        """初始化 DOCX 文档加载器。

        参数：
            path: 待解析的 .docx 文件路径。

        调用顺序：入库脚本或索引服务 -> PythonDocxLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """读取 DOCX 文件的段落和表格文本，返回 LangChain Document 列表。

        调用顺序：入库脚本或索引服务 -> PythonDocxLoader.load()。
        """
        if self.path.suffix.lower() != ".docx":
            raise RuntimeError(f"{self.path.suffix.lower()} 需要额外转换器，请先转为 .docx 后入库。")
        docx = DocxDocument(str(self.path))
        lines = [paragraph.text.strip() for paragraph in docx.paragraphs if paragraph.text.strip()]
        for table in docx.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        content = "\n".join(lines).strip()
        return [Document(page_content=content, metadata={"file_type": ".docx"})] if content else []


class PyMuPdfLoader:
    """轻量 PDF loader，按页提取文本层内容。

    调用顺序：入库脚本或索引服务 -> PyMuPdfLoader。
    """

    def __init__(self, path: Path) -> None:
        """初始化 PDF 文档加载器。

        参数：
            path: 待解析的 .pdf 文件路径。

        调用顺序：入库脚本或索引服务 -> PyMuPdfLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """按页提取 PDF 文本层内容，每页生成一个 LangChain Document。

        调用顺序：入库脚本或索引服务 -> PyMuPdfLoader.load()。
        """
        documents: list[Document] = []
        with fitz.open(str(self.path)) as pdf:
            for page_index, page in enumerate(pdf):
                text = page.get_text("text").strip()
                if text:
                    documents.append(Document(page_content=text, metadata={"page": page_index, "file_type": ".pdf"}))
        return documents


class PythonPptxLoader:
    """轻量 PPTX loader，按幻灯片提取文本框和表格文本。

    调用顺序：入库脚本或索引服务 -> PythonPptxLoader。
    """

    def __init__(self, path: Path) -> None:
        """初始化 PPTX 文档加载器。

        参数：
            path: 待解析的 .pptx 文件路径。

        调用顺序：入库脚本或索引服务 -> PythonPptxLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """按幻灯片提取文本框和表格文本，每页生成一个 LangChain Document。

        调用顺序：入库脚本或索引服务 -> PythonPptxLoader.load()。
        """
        if self.path.suffix.lower() != ".pptx":
            raise RuntimeError(f"{self.path.suffix.lower()} 需要额外转换器，请先转为 .pptx 后入库。")
        presentation = Presentation(str(self.path))
        documents: list[Document] = []
        for slide_index, slide in enumerate(presentation.slides):
            lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
                text = getattr(shape, "text", "").strip()
                if text:
                    lines.append(text)
            content = "\n".join(dict.fromkeys(lines)).strip()
            if content:
                documents.append(Document(page_content=content, metadata={"page": slide_index, "file_type": ".pptx"}))
        return documents


class ProjectTextLoader:
    """项目文本 loader，读取 UTF-8 文本并识别已复核 OCR Markdown。

    调用顺序：入库脚本或索引服务 -> ProjectTextLoader。
    """

    def __init__(self, path: Path) -> None:
        """初始化文本文件加载器。

        参数：
            path: 待读取的文本或 Markdown 文件路径。

        调用顺序：入库脚本或索引服务 -> ProjectTextLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """读取 UTF-8 文本文件内容，自动识别 OCR 复核元数据并返回 Document 列表。

        调用顺序：入库脚本或索引服务 -> ProjectTextLoader.load()。
        """
        content = self.path.read_text(encoding="utf-8")
        metadata = {"file_type": self.path.suffix.lower(), **parse_ocr_review_metadata(content)}
        return [Document(page_content=content, metadata=metadata)] if content.strip() else []


class TableDocumentLoader:
    """把 CSV/Excel 表格包装成 LangChain loader 协议。表格行转换成 Document 后复用主链路。

    调用顺序：入库脚本或索引服务 -> TableDocumentLoader。
    """

    def __init__(self, path: Path) -> None:
        """保存待加载的表格路径。

        调用顺序：入库脚本或索引服务 -> TableDocumentLoader.__init__()。
        """
        self.path = path

    def load(self) -> list[Document]:
        """加载表格并返回行级 Document。

        调用顺序：入库脚本或索引服务 -> TableDocumentLoader.load()。
        """
        return load_table_file(self.path)

def _table_loader(path: Path) -> SupportsLoad:
    """创建表格资料加载器。

    调用顺序：入库脚本或索引服务 -> _table_loader()。
    """
    return TableDocumentLoader(path)


DOCUMENT_LOADER_SPECS: tuple[DocumentLoaderSpec, ...] = (
    DocumentLoaderSpec(
        suffixes=(".txt", ".md"),
        factory=_utf8_text_loader,
        description="UTF-8 文本和 Markdown；Markdown 保留原文给标题切分器处理。",
    ),
    DocumentLoaderSpec(
        suffixes=(".pdf",),
        factory=_pdf_loader,
        description="PDF 文本层解析；扫描件 OCR 不默认进入主链路。",
    ),
    DocumentLoaderSpec(
        suffixes=(".docx", ".doc"),
        factory=_word_loader,
        description="Word 文档文本解析。",
    ),
    DocumentLoaderSpec(
        suffixes=(".ppt", ".pptx"),
        factory=_powerpoint_loader,
        description="PowerPoint 文本解析。",
    ),
    DocumentLoaderSpec(
        suffixes=(".html", ".htm"),
        factory=_docling_loader,
        description="HTML 资料解析；需要 DOCUMENT_PARSER_BACKEND=docling。",
    ),
    DocumentLoaderSpec(
        suffixes=(".csv", ".xlsx", ".xls"),
        factory=_table_loader,
        description="CSV/Excel 表格解析；按行保留表头、sheet 和单元格键值。",
    ),
)

DOCUMENT_LOADER_REGISTRY: dict[str, DocumentLoaderSpec] = {
    suffix: spec
    for spec in DOCUMENT_LOADER_SPECS
    for suffix in spec.suffixes
}
SUPPORTED_DOCUMENT_SUFFIXES = tuple(sorted(DOCUMENT_LOADER_REGISTRY))

def get_document_loader_spec(path: Path) -> DocumentLoaderSpec | None:
    """根据文件后缀获取加载器注册项。后缀标准化集中处理，避免散落。

    调用顺序：入库脚本或索引服务 -> get_document_loader_spec()。
    """
    return DOCUMENT_LOADER_REGISTRY.get(path.suffix.lower())

def load_file(path: Path) -> list[Document]:
    """把一个受支持的本地文件加载为 LangChain Document 对象。返回值统一为 Document，后续不关心原文件格式。

    调用顺序：入库脚本或索引服务 -> load_file()。
    """
    spec = get_document_loader_spec(path)
    if spec is None:
        raise ValueError(f"不支持的文档类型：{path}")
    if _use_docling_for(path):
        logger.debug("Loading document with Docling enhanced parser: %s", path)
        return _docling_loader(path).load()
    logger.debug("Loading document with %s: %s", spec.description, path)
    loader = spec.create_loader(path)
    return loader.load()

