"""Image and mixed-layout document risk detection for ingestion quality reports.

The online RAG chain must not assume that embedded images were indexed. This
module makes image-bearing files visible during offline ingestion quality
checks so operators can route scans and business-critical screenshots through
OCR review before activation.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


IMAGE_FILE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}
IMAGE_CONTAINER_SUFFIXES = {".pdf", ".docx", ".pptx"}
IMAGE_RISK_SUFFIXES = IMAGE_FILE_SUFFIXES | IMAGE_CONTAINER_SUFFIXES
BLOCKING_TEXT_CHAR_THRESHOLD = 120


@dataclass(frozen=True)
class ImageRisk:
    """A visible image-ingestion risk attached to a candidate source file."""

    path: str
    suffix: str
    severity: str
    image_count: int
    text_char_count: int
    reason: str
    requires_ocr_review: bool

    def as_dict(self) -> dict[str, Any]:
        """转换为可 JSON 序列化的诊断数据。

        调用顺序：测试或业务入口 -> ImageRisk.as_dict()。
        """
        return asdict(self)


def analyze_image_risk(path: Path, extracted_text: str = "") -> ImageRisk | None:
    """Detect image-bearing files that need OCR review or operator awareness."""

    suffix = path.suffix.lower()
    if suffix not in IMAGE_RISK_SUFFIXES:
        return None

    text_char_count = _visible_text_length(extracted_text)
    if suffix in IMAGE_FILE_SUFFIXES:
        return ImageRisk(
            path=str(path),
            suffix=suffix,
            severity="block",
            image_count=1,
            text_char_count=0,
            reason="独立图片不能直接进入知识库；必须先离线 OCR，人工复核后再提升为 Markdown 资料。",
            requires_ocr_review=True,
        )

    image_count = _embedded_image_count(path)
    if image_count <= 0:
        return None

    if text_char_count < BLOCKING_TEXT_CHAR_THRESHOLD:
        return ImageRisk(
            path=str(path),
            suffix=suffix,
            severity="block",
            image_count=image_count,
            text_char_count=text_char_count,
            reason="文件主要依赖图片或扫描页，文本层不足，必须先走 OCR 复核流程。",
            requires_ocr_review=True,
        )

    return ImageRisk(
        path=str(path),
        suffix=suffix,
        severity="review",
        image_count=image_count,
        text_char_count=text_char_count,
        reason="文件包含图片；当前入库只保证文本层进入知识库，若图片承载业务信息需单独 OCR/复核。",
        requires_ocr_review=False,
    )


def _visible_text_length(text: str) -> int:
    return len("".join(str(text or "").split()))


def _embedded_image_count(path: Path) -> int:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _pdf_image_count(path)
    if suffix == ".docx":
        return _docx_image_count(path)
    if suffix == ".pptx":
        return _pptx_image_count(path)
    return 0


def _pdf_image_count(path: Path) -> int:
    with fitz.open(str(path)) as document:
        return sum(len(page.get_images(full=True)) for page in document)


def _docx_image_count(path: Path) -> int:
    document = DocxDocument(str(path))
    inline_count = len(document.inline_shapes)
    relation_count = sum(1 for relation in document.part.rels.values() if "image" in relation.reltype)
    return max(inline_count, relation_count)


def _pptx_image_count(path: Path) -> int:
    presentation = Presentation(str(path))
    return sum(_count_shape_images(slide.shapes) for slide in presentation.slides)


def _count_shape_images(shapes: Any) -> int:
    count = 0
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            count += 1
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            count += _count_shape_images(nested)
    return count
